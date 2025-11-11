from werkzeug.security import generate_password_hash, check_password_hash
from flask import Flask, render_template, request, redirect, url_for, flash, session, make_response, abort, send_file
import os
import psycopg2
from psycopg2.extras import RealDictCursor
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import Color
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from openpyxl import Workbook
from openpyxl.styles import Alignment

app = Flask(__name__)
app.secret_key = "ugel_lauricocha_2025"

# Conexión a PostgreSQL
def get_db_connection():
    conn = psycopg2.connect(
        os.environ['DATABASE_URL'],
        cursor_factory=RealDictCursor
    )
    return conn

# Inicializar base de datos
def init_db():
    conn = get_db_connection()
    cur = conn.cursor()
    # Tabla de usuarios
    cur.execute('''
        CREATE TABLE IF NOT EXISTS usuarios (
            id SERIAL PRIMARY KEY,
            usuario TEXT UNIQUE NOT NULL,
            contrasena TEXT NOT NULL,
            nombre_completo TEXT,
            rol TEXT NOT NULL CHECK (rol IN ('especialista', 'jefe', 'admin')),
            activo BOOLEAN DEFAULT TRUE,
            creado_en TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    # Tabla de visitas
    cur.execute('''
        CREATE TABLE IF NOT EXISTS visitas (
            id SERIAL PRIMARY KEY,
            usuario_id INTEGER REFERENCES usuarios(id),
            numero_informe TEXT UNIQUE,
            fecha TEXT,
            institucion TEXT,
            nivel TEXT,
            tipo_visita TEXT,
            fortalezas TEXT,
            mejoras TEXT,
            recomendaciones TEXT,
            compromisos TEXT,
            creado_en TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    # Usuario administrador (con contraseña hasheada)
    admin_pass = generate_password_hash('123456')
    cur.execute('''
        INSERT INTO usuarios (usuario, contrasena, nombre_completo, rol, activo)
        VALUES (%s, %s, %s, %s, %s)
        ON CONFLICT (usuario) DO NOTHING
    ''', ('admin', admin_pass, 'Administrador', 'admin', True))
    
    conn.commit()
    cur.close()
    conn.close()

@app.route('/')
def index():
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        usuario = request.form['usuario']
        contrasena = request.form['contrasena']
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute("SELECT * FROM usuarios WHERE usuario = %s", (usuario,))
        user = cur.fetchone()
        conn.close()
        
        if user and user['activo'] and check_password_hash(user['contrasena'], contrasena):
            session['user_id'] = user['id']
            session['rol'] = user['rol']
            session['nombre'] = user['nombre_completo']
            return redirect(url_for('dashboard'))
        else:
            flash('Usuario inactivo o credenciales incorrectos')
    
    # Respuesta sin caché
    resp = make_response(render_template('login.html'))
    resp.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
    resp.headers["Pragma"] = "no-cache"
    resp.headers["Expires"] = "0"
    return resp

@app.route('/dashboard')
def dashboard():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    user_id = session['user_id']
    rol = session['rol']
    mes_filtro = request.args.get('mes', '')

    conn = get_db_connection()
    cur = conn.cursor()
    
    # Construir condiciones de filtrado
    where_conditions = []
    params = []
    
    if mes_filtro:
        where_conditions.append("fecha LIKE %s")
        params.append(f"{mes_filtro}%")
    
    if rol not in ['admin', 'jefe']:
        where_conditions.append("usuario_id = %s")
        params.append(user_id)
    
    where_clause = " WHERE " + " AND ".join(where_conditions) if where_conditions else ""
    query = "SELECT * FROM visitas" + where_clause + " ORDER BY id DESC"
    cur.execute(query, params)
    visitas = cur.fetchall()
    
    # Totales
    total_visitas = len(visitas)
    visitas_mes = 0
    porcentaje_meta = 0

    if mes_filtro:
        visitas_mes = total_visitas
    else:
        mes_actual = datetime.now().strftime("%Y-%m")
        count_params = [f"{mes_actual}%"]
        if rol not in ['admin', 'jefe']:
            cur.execute("SELECT COUNT(*) FROM visitas WHERE fecha LIKE %s AND usuario_id = %s", count_params + [user_id])
        else:
            cur.execute("SELECT COUNT(*) FROM visitas WHERE fecha LIKE %s", (f"{mes_actual}%",))
        visitas_mes = cur.fetchone()['count']
    
    meta_mensual = 30
    porcentaje_meta = (visitas_mes / meta_mensual * 100) if meta_mensual > 0 else 0

    # Estadísticas
    nivel_params = params.copy()
    tipo_params = params.copy()
    nivel_query = "SELECT nivel, COUNT(*) FROM visitas" + where_clause + " GROUP BY nivel"
    tipo_query = "SELECT tipo_visita, COUNT(*) FROM visitas" + where_clause + " GROUP BY tipo_visita"

    cur.execute(nivel_query, nivel_params)
    niveles = dict(cur.fetchall())
    
    cur.execute(tipo_query, tipo_params)
    tipos = dict(cur.fetchall())
    
    conn.close()
    
    return render_template('dashboard.html', 
                         visitas=visitas,
                         total_visitas=total_visitas,
                         visitas_mes=visitas_mes,
                         porcentaje_meta=porcentaje_meta,
                         niveles=niveles,
                         tipos=tipos,
                         mes_filtro=mes_filtro)

def generar_numero_informe():
    """Genera un número de informe único: INF-2025-001, etc., usando el último ID."""
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT COALESCE(MAX(id), 0) FROM visitas")
    last_id = cur.fetchone()['coalesce'] + 1
    conn.close()
    año = datetime.now().year
    return f"INF-{año}-{last_id:03d}"

@app.route('/nueva_visita', methods=['GET', 'POST'])
def nueva_visita():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        user_id = session['user_id']
        fecha = request.form['fecha']
        institucion = request.form['institucion']
        nivel = request.form['nivel']
        tipo = request.form['tipo']
        fortalezas = request.form['fortalezas']
        mejoras = request.form['mejoras']
        recomendaciones = request.form['recomendaciones']
        compromisos = request.form.get('compromisos', '')

        conn = get_db_connection()
        cur = conn.cursor()
        numero_informe = generar_numero_informe()
        cur.execute('''
            INSERT INTO visitas (
                usuario_id, numero_informe, fecha, institucion, nivel, tipo_visita,
                fortalezas, mejoras, recomendaciones, compromisos
            ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        ''', (
            user_id, numero_informe, fecha, institucion, nivel, tipo,
            fortalezas, mejoras, recomendaciones, compromisos
        ))
        conn.commit()
        conn.close()
        flash('Visita registrada con éxito')
        return redirect(url_for('dashboard'))
    return render_template('nueva_visita.html')

# === Rutas de generación de informes (PPTX, PDF, Excel) ===

@app.route('/generar_ppt/<int:visita_id>')
def generar_ppt(visita_id):
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    user_id = session['user_id']
    rol = session['rol']

    conn = get_db_connection()
    cur = conn.cursor()

    if rol in ['admin', 'jefe']:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s
        """, (visita_id,))
    else:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s AND v.usuario_id = %s
        """, (visita_id, user_id))
    
    visita = cur.fetchone()
    conn.close()

    if not visita:
        flash("Visita no encontrada o no autorizada")
        return redirect(url_for('dashboard'))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logo_path = os.path.join(os.path.dirname(__file__), 'templates_ppt', 'logo.png')
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), height=Inches(0.8))

    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"REPORTE DE VISITA PEDAGÓGICA\n{visita['numero_informe']}"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(10), Inches(1))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = f"Institución: {visita['institucion']}\nEspecialista: {visita['especialista_nombre']}\nFecha: {visita['fecha']} | Tipo: {visita['tipo_visita']}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.alignment = PP_ALIGN.CENTER

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf3 = title2_box.text_frame
    tf3.clear()
    p3 = tf3.paragraphs[0]
    p3.text = "OBSERVACIONES ESTRUCTURADAS"
    p3.font.bold = True
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0, 51, 102)

    obs_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf4 = obs_box.text_frame
    tf4.clear()
    p4 = tf4.paragraphs[0]
    p4.text = f"""✅ Fortalezas:\n{visita['fortalezas'] or 'No registradas.'}

⚠️ Áreas de mejora:\n{visita['mejoras'] or 'No registradas.'}

💡 Recomendaciones:\n{visita['recomendaciones'] or 'No registradas.'}

📝 Compromisos:\n{visita['compromisos'] or 'No registrados.'}"""
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0, 0, 0)

    filename = f"reporte_visita_{visita_id}.pptx"
    filepath = os.path.join(os.path.dirname(__file__), filename)
    prs.save(filepath)

    return send_file(filepath, as_attachment=True)

@app.route('/generar_pdf/<int:visita_id>')
def generar_pdf(visita_id):
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    user_id = session['user_id']
    rol = session['rol']

    conn = get_db_connection()
    cur = conn.cursor()

    if rol in ['admin', 'jefe']:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s
        """, (visita_id,))
    else:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s AND v.usuario_id = %s
        """, (visita_id, user_id))
    
    visita = cur.fetchone()
    conn.close()
    
    if not visita:
        flash("Visita no encontrada o no autorizada")
        return redirect(url_for('dashboard'))

    filename = f"informe_visita_{visita_id}.pdf"
    filepath = os.path.join(os.path.dirname(__file__), filename)

    doc = SimpleDocTemplate(filepath, pagesize=letter,
                            rightMargin=72, leftMargin=72,
                            topMargin=72, bottomMargin=72)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=16,
        spaceAfter=14,
        textColor=Color(0, 51/255, 102/255),
        alignment=TA_CENTER
    )

    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=10,
        alignment=TA_LEFT
    )

    story = []

    logo_path = os.path.join(os.path.dirname(__file__), 'templates_ppt', 'logo.png')
    if os.path.exists(logo_path):
        im = Image(logo_path, 1.5*inch, 0.8*inch)
        im.hAlign = 'CENTER'
        story.append(im)
        story.append(Spacer(1, 24))

    story.append(Paragraph("INFORME DE VISITA PEDAGÓGICA", title_style))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Número de Informe:</b> {visita['numero_informe']}", normal_style))
    story.append(Spacer(1, 12))

    datos = [
        f"<b>Institución Educativa:</b> {visita['institucion']}",
        f"<b>Nivel:</b> {visita['nivel']}",
        f"<b>Tipo de Visita:</b> {visita['tipo_visita']}",
        f"<b>Especialista:</b> {visita['especialista_nombre']}",
        f"<b>Fecha:</b> {visita['fecha']}",
    ]

    for dato in datos:
        story.append(Paragraph(dato, normal_style))

    story.append(Spacer(1, 18))
    story.append(Paragraph("<b>OBSERVACIONES ESTRUCTURADAS</b>", normal_style))
    story.append(Spacer(1, 6))

    story.append(Paragraph("<b>✅ Fortalezas:</b>", normal_style))
    story.append(Paragraph(visita['fortalezas'] or "No registradas.", normal_style))
    story.append(Spacer(1, 6))

    story.append(Paragraph("<b>⚠️ Áreas de mejora:</b>", normal_style))
    story.append(Paragraph(visita['mejoras'] or "No registradas.", normal_style))
    story.append(Spacer(1, 6))

    story.append(Paragraph("<b>💡 Recomendaciones:</b>", normal_style))
    story.append(Paragraph(visita['recomendaciones'] or "No registradas.", normal_style))
    story.append(Spacer(1, 6))

    story.append(Paragraph("<b>📝 Compromisos del docente:</b>", normal_style))
    story.append(Paragraph(visita['compromisos'] or "No registrados.", normal_style))

    story.append(Spacer(1, 36))
    story.append(Paragraph("_____________________________________", normal_style))
    story.append(Paragraph("Firma del Especialista", normal_style))
    story.append(Spacer(1, 12))
    story.append(Paragraph("_____________________________________", normal_style))
    story.append(Paragraph("Sello de la Institución Educativa", normal_style))

    doc.build(story)
    return send_file(filepath, as_attachment=True)

@app.route('/gestion_usuarios')
def gestion_usuarios():
    if 'rol' not in session or session['rol'] != 'admin':
        flash('Acceso no autorizado')
        return redirect(url_for('dashboard'))
    
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT id, usuario, nombre_completo, rol, activo, creado_en FROM usuarios ORDER BY creado_en DESC")
    usuarios = cur.fetchall()
    conn.close()
    return render_template('gestion_usuarios.html', usuarios=usuarios)

@app.route('/crear_usuario', methods=['POST'])
def crear_usuario():
    if 'rol' not in session or session['rol'] != 'admin':
        flash('Acceso no autorizado')
        return redirect(url_for('dashboard'))
    
    usuario = request.form['usuario']
    contrasena = request.form['contrasena']
    nombre_completo = request.form['nombre_completo']
    rol = request.form['rol']
    
    contrasena_hash = generate_password_hash(contrasena)
    
    conn = get_db_connection()
    cur = conn.cursor()
    try:
        cur.execute('''
            INSERT INTO usuarios (usuario, contrasena, nombre_completo, rol, activo)
            VALUES (%s, %s, %s, %s, %s)
        ''', (usuario, contrasena_hash, nombre_completo, rol, True))
        conn.commit()
        flash('Usuario creado exitosamente')
    except psycopg2.IntegrityError:
        conn.rollback()
        flash('Error: El nombre de usuario ya existe')
    finally:
        conn.close()
    
    return redirect(url_for('gestion_usuarios'))

@app.route('/exportar_excel')
def exportar_excel():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    user_id = session['user_id']
    rol = session['rol']
    conn = get_db_connection()
    cur = conn.cursor()

    if rol in ['admin', 'jefe']:
        cur.execute("""
            SELECT 
                v.numero_informe, v.fecha, v.institucion, v.nivel, v.tipo_visita,
                u.nombre_completo as especialista,
                v.fortalezas, v.mejoras, v.recomendaciones, v.compromisos
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            ORDER BY v.id DESC
        """)
    else:
        cur.execute("""
            SELECT 
                v.numero_informe, v.fecha, v.institucion, v.nivel, v.tipo_visita,
                u.nombre_completo as especialista,
                v.fortalezas, v.mejoras, v.recomendaciones, v.compromisos
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.usuario_id = %s
            ORDER BY v.id DESC
        """, (user_id,))
    
    visitas = cur.fetchall()
    conn.close()

    wb = Workbook()
    ws = wb.active
    ws.title = "Visitas Pedagógicas"

    headers = [
        "N° Informe", "Fecha", "Institución", "Nivel", "Tipo de Visita", "Especialista",
        "Fortalezas", "Áreas de mejora", "Recomendaciones", "Compromisos"
    ]
    ws.append(headers)

    for visita in visitas:
        ws.append([
            visita['numero_informe'],
            visita['fecha'],
            visita['institucion'],
            visita['nivel'],
            visita['tipo_visita'],
            visita['especialista'],
            visita['fortalezas'] or "",
            visita['mejoras'] or "",
            visita['recomendaciones'] or "",
            visita['compromisos'] or ""
        ])

    for col in ws.columns:
        max_length = 0
        column = col[0].column_letter
        for cell in col:
            if cell.value is not None:
                max_length = max(max_length, len(str(cell.value)))
        adjusted_width = max(max_length + 2, 15)
        ws.column_dimensions[column].width = min(adjusted_width, 50)
        for cell in col:
            if cell.row > 1:
                cell.alignment = Alignment(wrap_text=True)

    filename = "visitas_pedagogicas.xlsx"
    filepath = os.path.join(os.path.dirname(__file__), filename)
    wb.save(filepath)

    return send_file(filepath, as_attachment=True)

@app.route('/editar_usuario/<int:usuario_id>')
def editar_usuario(usuario_id):
    if 'rol' not in session or session['rol'] != 'admin':
        flash('Acceso no autorizado')
        return redirect(url_for('dashboard'))
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT * FROM usuarios WHERE id = %s", (usuario_id,))
    usuario = cur.fetchone()
    conn.close()
    if not usuario:
        flash('Usuario no encontrado')
        return redirect(url_for('gestion_usuarios'))
    return render_template('editar_usuario.html', usuario=usuario)

@app.route('/actualizar_usuario/<int:usuario_id>', methods=['POST'])
def actualizar_usuario(usuario_id):
    if 'rol' not in session or session['rol'] != 'admin':
        flash('Acceso no autorizado')
        return redirect(url_for('dashboard'))
    
    nombre_completo = request.form['nombre_completo']
    rol = request.form['rol']
    activo = 'activo' in request.form
    contrasena = request.form['contrasena']
    
    conn = get_db_connection()
    cur = conn.cursor()
    
    if contrasena:
        contrasena_hash = generate_password_hash(contrasena)
        cur.execute('''
            UPDATE usuarios 
            SET nombre_completo = %s, rol = %s, activo = %s, contrasena = %s
            WHERE id = %s
        ''', (nombre_completo, rol, activo, contrasena_hash, usuario_id))
    else:
        cur.execute('''
            UPDATE usuarios 
            SET nombre_completo = %s, rol = %s, activo = %s
            WHERE id = %s
        ''', (nombre_completo, rol, activo, usuario_id))
    
    conn.commit()
    conn.close()
    flash('Usuario actualizado exitosamente')
    return redirect(url_for('gestion_usuarios'))   

@app.route('/eliminar_usuario/<int:usuario_id>', methods=['POST'])
def eliminar_usuario(usuario_id):
    if 'rol' not in session or session['rol'] != 'admin':
        flash('Acceso no autorizado')
        return redirect(url_for('dashboard'))
    
    if usuario_id == session['user_id']:
        flash('No puedes eliminarte a ti mismo')
        return redirect(url_for('gestion_usuarios'))
    
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("UPDATE usuarios SET activo = FALSE WHERE id = %s", (usuario_id,))
    conn.commit()
    conn.close()
    flash('Usuario desactivado exitosamente')
    return redirect(url_for('gestion_usuarios'))

@app.route('/generar_informe_mensual/<mes>')
def generar_informe_mensual(mes):
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cur = conn.cursor()
    
    if mes == 'todos':
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            ORDER BY v.fecha DESC
        """)
    else:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v
            JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.fecha LIKE %s
            ORDER BY v.fecha DESC
        """, (f"{mes}%",))
    
    visitas = cur.fetchall()
    
    total = len(visitas)
    niveles = {}
    tipos = {}
    
    if total > 0:
        if mes == 'todos':
            cur.execute("SELECT nivel, COUNT(*) FROM visitas GROUP BY nivel")
            niveles = {row['nivel']: row['count'] for row in cur.fetchall()}
            cur.execute("SELECT tipo_visita, COUNT(*) FROM visitas GROUP BY tipo_visita")
            tipos = {row['tipo_visita']: row['count'] for row in cur.fetchall()}
        else:
            cur.execute("SELECT nivel, COUNT(*) FROM visitas WHERE fecha LIKE %s GROUP BY nivel", (f"{mes}%",))
            niveles = {row['nivel']: row['count'] for row in cur.fetchall()}
            cur.execute("SELECT tipo_visita, COUNT(*) FROM visitas WHERE fecha LIKE %s GROUP BY tipo_visita", (f"{mes}%",))
            tipos = {row['tipo_visita']: row['count'] for row in cur.fetchall()}
    
    conn.close()
    
    return generar_pdf_informe_mensual(visitas, mes, total, niveles, tipos)

def generar_pdf_informe_mensual(visitas, mes, total, niveles, tipos):
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    
    filename = f"informe_mensual_{mes}.pdf"
    filepath = os.path.join(os.path.dirname(__file__), filename)
    
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    logo_path = os.path.join(os.path.dirname(__file__), 'templates_ppt', 'logo.png')
    if os.path.exists(logo_path):
        im = Image(logo_path, 1.5*inch, 0.8*inch)
        im.hAlign = 'CENTER'
        story.append(im)
        story.append(Spacer(1, 24))
    
    story.append(Paragraph("<b>UNIDAD DE GESTIÓN EDUCATIVA LOCAL LAURICOCHA</b>", styles['Title']))
    story.append(Paragraph("<b>ÁREA DE GESTIÓN PEDAGÓGICA</b>", styles['Heading2']))
    story.append(Spacer(1, 12))
    
    mes_str = "TODOS LOS MESES" if mes == 'todos' else mes
    story.append(Paragraph(f"<b>INFORME MENSUAL DE MONITOREO PEDAGÓGICO - {mes_str}</b>", styles['Heading1']))
    story.append(Spacer(1, 12))
    
    story.append(Paragraph("<b>RESUMEN ESTADÍSTICO</b>", styles['Heading3']))
    story.append(Paragraph(f"Total de visitas realizadas: {total}", styles['Normal']))
    
    if niveles:
        niveles_str = ", ".join([f"{k}: {v}" for k, v in niveles.items()])
        story.append(Paragraph(f"Visitas por nivel: {niveles_str}", styles['Normal']))
    
    if tipos:
        tipos_str = ", ".join([f"{k}: {v}" for k, v in tipos.items()])
        story.append(Paragraph(f"Visitas por tipo: {tipos_str}", styles['Normal']))
    
    story.append(Spacer(1, 18))
    
    if visitas:
        story.append(Paragraph("<b>DETALLE DE VISITAS</b>", styles['Heading3']))
        data = [["N° Informe", "Fecha", "Institución", "Nivel", "Tipo", "Especialista"]]
        for v in visitas:
            data.append([
                Paragraph(v['numero_informe'], styles['Normal']),
                Paragraph(v['fecha'], styles['Normal']),
                Paragraph(v['institucion'], styles['Normal']),
                Paragraph(v['nivel'], styles['Normal']),
                Paragraph(v['tipo_visita'], styles['Normal']),
                Paragraph(v['especialista_nombre'], styles['Normal'])
            ])
        
        col_widths = [60, 80, 100, 50, 70, 80]
        table = Table(data, colWidths=col_widths)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.gray),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE', (0,0), (-1,0), 8),
            ('BOTTOMPADDING', (0,0), (-1,0), 12),
            ('GRID', (0,0), (-1,-1), 1, colors.black),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ]))
        story.append(table)
    
    story.append(Spacer(1, 24))
    story.append(Paragraph("<b>OBSERVACIONES ESTRUCTURADAS COMPLETAS</b>", styles['Heading3']))
    
    for v in visitas:
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"<b>{v['numero_informe']} - {v['institucion']}</b>", styles['Heading4']))
        story.append(Paragraph("<b>✅ Fortalezas:</b>", styles['Normal']))
        story.append(Paragraph(v['fortalezas'] or "No registradas.", styles['Normal']))
        story.append(Spacer(1, 6))
        story.append(Paragraph("<b>⚠️ Áreas de mejora:</b>", styles['Normal']))
        story.append(Paragraph(v['mejoras'] or "No registradas.", styles['Normal']))
        story.append(Spacer(1, 6))
        story.append(Paragraph("<b>💡 Recomendaciones:</b>", styles['Normal']))
        story.append(Paragraph(v['recomendaciones'] or "No registradas.", styles['Normal']))
        story.append(Spacer(1, 6))
        story.append(Paragraph("<b>📝 Compromisos del docente:</b>", styles['Normal']))
        story.append(Paragraph(v['compromisos'] or "No registrados.", styles['Normal']))
    
    story.append(Spacer(1, 36))
    story.append(Paragraph("_____________________________________", styles['Normal']))
    story.append(Paragraph("Firma del Jefe del Área de Gestión Pedagógica", styles['Normal']))
    story.append(Paragraph("UGEL Lauricocha", styles['Normal']))
    
    doc.build(story)
    return send_file(filepath, as_attachment=True)

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/actualizar_visitas')
def actualizar_visitas():
    if session.get('rol') != 'admin':
        return "Acceso denegado", 403
    conn = get_db_connection()
    cur = conn.cursor()
    try:
        # Verificar si ya existen las columnas
        cur.execute("SELECT column_name FROM information_schema.columns WHERE table_name = 'visitas'")
        cols = [row['column_name'] for row in cur.fetchall()]
        cambios = []
        if 'fortalezas' not in cols:
            cur.execute("ALTER TABLE visitas ADD COLUMN fortalezas TEXT")
            cambios.append('fortalezas')
        if 'mejoras' not in cols:
            cur.execute("ALTER TABLE visitas ADD COLUMN mejoras TEXT")
            cambios.append('mejoras')
        if 'recomendaciones' not in cols:
            cur.execute("ALTER TABLE visitas ADD COLUMN recomendaciones TEXT")
            cambios.append('recomendaciones')
        if 'compromisos' not in cols:
            cur.execute("ALTER TABLE visitas ADD COLUMN compromisos TEXT")
            cambios.append('compromisos')
        conn.commit()
        if cambios:
            return f"✅ Columnas añadidas: {', '.join(cambios)}"
        else:
            return "✅ Todas las columnas ya existen"
    except Exception as e:
        return f"⚠️ Error: {str(e)}"
    finally:
        conn.close()

@app.route('/verificar_columnas')
def verificar_columnas():
    if session.get('rol') != 'admin':
        return "Acceso denegado", 403
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT column_name FROM information_schema.columns WHERE table_name = 'visitas'")
    columnas = [row['column_name'] for row in cur.fetchall()]
    conn.close()
    return f"Columnas en 'visitas': {', '.join(columnas)}"

@app.route('/recursos')
def recursos():
    try:
        if 'user_id' not in session:
            return redirect(url_for('login'))
        
        recursos_path = os.path.join(os.path.dirname(__file__), 'recursos')
        if not os.path.exists(recursos_path):
            return "<h2>Error: Carpeta 'recursos' no encontrada en el servidor.</h2>", 500
        
        recursos = {
            'rúbricas': [
                {'nombre': 'Rúbrica de Evaluación - Inicial', 'descripcion': 'Instrumento para evaluar competencias en comunicación y matemática.', 'archivo': 'rubrica_inicial.pdf'},
                {'nombre': 'Rúbrica de Evaluación - Primaria', 'descripcion': 'Evaluación de desempeños en áreas curriculares según el CNEB.', 'archivo': 'rubrica_primaria.pdf'},
                {'nombre': 'Rúbrica de Evaluación - Secundaria', 'descripcion': 'Enfoque en competencias ciudadanas y científico-tecnológicas.', 'archivo': 'rubrica_secundaria.pdf'}
            ],
            'guías': [
                {'nombre': 'Guía de Monitoreo Pedagógico 2025', 'descripcion': 'Protocolo oficial del MINEDU para visitas de monitoreo y asesoría.', 'archivo': 'guia_monitoreo_2025.pdf'},
                {'nombre': 'Protocolo de Observación de Aula', 'descripcion': 'Registro estructurado de prácticas pedagógicas efectivas.', 'archivo': 'protocolo_observacion.pdf'}
            ],
            'planes': [
                {'nombre': 'Modelo de Plan de Mejora Institucional', 'descripcion': 'Plantilla editable para instituciones educativas (Word).', 'archivo': 'plan_mejora_modelo.docx'},
                {'nombre': 'Formato de Compromisos de Gestión Escolar', 'descripcion': 'Seguimiento a metas del Proyecto Educativo Institucional.', 'archivo': 'compromisos_gestion.pdf'}
            ],
            'normas': [
                {'nombre': 'Directiva UGEL Lauricocha N° 001-2025', 'descripcion': 'Lineamientos para el monitoreo pedagógico en zonas rurales.', 'archivo': 'directiva_ugel_lauricocha.pdf'},
                {'nombre': 'Resolución Ministerial N° 123-2025-MINEDU', 'descripcion': 'Normas para la evaluación del aprendizaje en educación básica.', 'archivo': 'resolucion_minedu.pdf'}
            ]
        }
        return render_template('recursos.html', recursos=recursos)
    except Exception as e:
        return f"<h2>Error en el repositorio: {str(e)}</h2>", 500

@app.route('/recursos/<categoria>/<archivo>')
def descargar_recurso(categoria, archivo):
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    categorias_validas = ['rúbricas', 'guías', 'planes', 'normas']
    if categoria not in categorias_validas:
        abort(404)
    
    try:
        return send_file(os.path.join('recursos', categoria, archivo), as_attachment=True)
    except FileNotFoundError:
        abort(404)

@app.route('/migrar_activo_2025')
def migrar_activo_2025():
    # Protección: solo si hay sesión y es admin
    if 'user_id' not in session or session.get('rol') != 'admin':
        return "Acceso denegado. Inicia sesión como admin.", 403
    
    conn = get_db_connection()
    cur = conn.cursor()
    try:
        # Añadir columna 'activo' si no existe
        cur.execute("ALTER TABLE usuarios ADD COLUMN IF NOT EXISTS activo BOOLEAN DEFAULT TRUE")
        # Asegurar que todos los usuarios existentes estén activos
        cur.execute("UPDATE usuarios SET activo = TRUE WHERE activo IS NULL OR activo != TRUE")
        conn.commit()
        return """
        <h2>✅ ¡Migración completada!</h2>
        <p>Columna 'activo' añadida y todos los usuarios activados.</p>
        <p><b>Recomendación:</b> Elimina esta ruta del código y haz un nuevo push a Git.</p>
        """
    except Exception as e:
        conn.rollback()
        return f"<h2>❌ Error:</h2><pre>{str(e)}</pre>"
    finally:
        conn.close()
        
# --- Inicialización ---
if __name__ == '__main__':
    init_db()  # Asegura que la DB esté lista al iniciar
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port)