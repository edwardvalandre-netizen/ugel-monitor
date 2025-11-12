from flask import session, redirect, url_for, flash, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from utils.db import get_db_connection
import os

def generar_ppt_visita(visita_id):
    if 'user_id' not in session:
        return redirect(url_for('auth.login'))

    user_id = session['user_id']
    rol = session['rol']

    conn = get_db_connection()
    cur = conn.cursor()

    if rol in ['admin', 'jefe']:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s
        """, (visita_id,))
    else:
        cur.execute("""
            SELECT v.*, u.nombre_completo as especialista_nombre
            FROM visitas v JOIN usuarios u ON v.usuario_id = u.id
            WHERE v.id = %s AND v.usuario_id = %s
        """, (visita_id, user_id))

    visita = cur.fetchone()
    conn.close()

    if not visita:
        flash("Visita no encontrada o no autorizada")
        return redirect(url_for('visitas.dashboard'))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logo_path = os.path.join(os.path.dirname(__file__), 'templates_ppt', 'logo.png')
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3), height=Inches(0.8))

    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"REPORTE DE VISITA PEDAGÓGICA\n{visita['numero_informe']}"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(10), Inches(1))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = f"Institución: {visita['institucion']}\nEspecialista: {visita['especialista_nombre']}\nFecha: {visita['fecha']} | Tipo: {visita['tipo_visita']}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.alignment = PP_ALIGN.CENTER

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    obs_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf3 = obs_box.text_frame
    tf3.clear()
    p3 = tf3.paragraphs[0]
    p3.text = f"""✅ Fortalezas:\n{visita['fortalezas'] or 'No registradas.'}

⚠️ Áreas de mejora:\n{visita['mejoras'] or 'No registradas.'}

💡 Recomendaciones:\n{visita['recomendaciones'] or 'No registradas.'}

📝 Compromisos:\n{visita['compromisos'] or 'No registrados.'}"""
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0, 0, 0)

    filename = f"reporte_visita_{visita_id}.pptx"
    filepath = os.path.join(os.path.dirname(__file__), filename)
    prs.save(filepath)

    return send_file(filepath, as_attachment=True)
